"""SharePoint / OneDrive — Read-Only local access via CloudStorage.
Searches, lists, and reads files synced to ~/Library/CloudStorage.
Supports txt/md/csv/json, pptx, docx, xlsx.
No Graph API, no OAuth — reads local OneDrive-synced files directly."""

import asyncio
import json
import logging
import os
import subprocess
from pathlib import Path
from typing import List, Optional

from tools.base import BaseTool, make_tool_def

logger = logging.getLogger("sharepoint-local")

# ---------------------------------------------------------------------------
# Safe roots — the ONLY directories we will read from
# ---------------------------------------------------------------------------

_CLOUD_BASE = Path.home() / "Library" / "CloudStorage"

# Auto-discover all OneDrive and SharePoint folders instead of hardcoding tenant names
_CANDIDATE_ROOTS = []
if _CLOUD_BASE.exists():
    _CANDIDATE_ROOTS = sorted([
        p for p in _CLOUD_BASE.iterdir()
        if p.is_dir() and (p.name.startswith("OneDrive") or "SharePoint" in p.name or "Salesforce" in p.name)
    ])

ALLOWED_ROOTS = [r for r in _CANDIDATE_ROOTS if r.exists()]


def _safe_path(path_str: str) -> Path:
    """Resolve path and verify it sits inside an allowed root."""
    p = Path(path_str).expanduser().resolve()
    for root in ALLOWED_ROOTS:
        try:
            p.relative_to(root)
            return p
        except ValueError:
            continue
    # Also accept paths already inside _CLOUD_BASE (catches new roots)
    try:
        p.relative_to(_CLOUD_BASE.resolve())
        return p
    except ValueError:
        pass
    raise ValueError(
        f"Path '{path_str}' is outside allowed SharePoint/OneDrive directories. "
        f"Only files under ~/Library/CloudStorage/ are accessible."
    )


def _run_search(query: str, file_types: str, root: str, max_results: int) -> str:
    """Search across synced SharePoint/OneDrive libraries using Spotlight."""
    max_results = min(max_results, 200)

    if root:
        try:
            scope_path = _safe_path(root)
        except ValueError as e:
            return json.dumps({"error": "Operation failed"})
        search_scopes = [str(scope_path)]
    else:
        search_scopes = [str(r) for r in ALLOWED_ROOTS]

    results = []

    for scope in search_scopes:
        if not os.path.isdir(scope):
            continue

        name_query = f'kMDItemFSName == "*{query}*"cd'
        content_query = f'kMDItemTextContent == "*{query}*"cd'
        combined = f'({name_query}) || ({content_query})'

        cmd = ["mdfind", "-onlyin", scope, combined]
        try:
            out = subprocess.run(cmd, capture_output=True, text=True, timeout=15)
            paths = [p.strip() for p in out.stdout.splitlines() if p.strip()]
        except subprocess.TimeoutExpired:
            cmd2 = ["find", scope, "-iname", f"*{query}*", "-not", "-name", ".*"]
            out2 = subprocess.run(cmd2, capture_output=True, text=True, timeout=15)
            paths = [p.strip() for p in out2.stdout.splitlines() if p.strip()]
        except Exception as e:
            logger.error(f"Search error in {scope}: {e}")
            continue

        if file_types:
            exts = set()
            for x in file_types.split(","):
                x = x.strip().lstrip(".")
                if x:
                    exts.add(f".{x}")
            paths = [p for p in paths if Path(p).suffix.lower() in exts]

        for path_str in paths:
            p = Path(path_str)
            if not p.exists() or p.name.startswith("."):
                continue
            try:
                stat = p.stat()
                file_type = "folder" if p.is_dir() else p.suffix.lower().lstrip(".")
                results.append({
                    "path": str(p),
                    "name": p.name,
                    "size_kb": round(stat.st_size / 1024, 1),
                    "modified": stat.st_mtime,
                    "type": file_type,
                })
            except OSError:
                continue

        if len(results) >= max_results:
            break

    results = results[:max_results]
    results.sort(key=lambda x: x.get("modified", 0), reverse=True)

    return json.dumps({
        "query": query,
        "count": len(results),
        "results": results,
    }, indent=2)


def _run_list(folder_path: str, depth: int) -> str:
    """List files and folders in a SharePoint/OneDrive directory."""
    depth = min(max(depth, 1), 3)

    if not folder_path:
        roots_info = []
        for r in ALLOWED_ROOTS:
            try:
                items = list(r.iterdir())
                roots_info.append({
                    "path": str(r),
                    "name": r.name,
                    "children_count": len(items),
                    "type": "library_root",
                })
            except PermissionError:
                pass
        return json.dumps({"roots": roots_info, "count": len(roots_info)}, indent=2)

    try:
        base = _safe_path(folder_path)
    except ValueError as e:
        return json.dumps({"error": "Operation failed"})

    if not base.exists():
        return json.dumps({"error": f"Path does not exist: {folder_path}"})

    def _list_dir(p: Path, current_depth: int) -> list:
        items = []
        try:
            entries = sorted(p.iterdir(), key=lambda x: (x.is_file(), x.name.lower()))
        except PermissionError:
            return []
        for entry in entries:
            if entry.name.startswith(".") or entry.name == "Icon":
                continue
            try:
                stat = entry.stat()
                item = {
                    "name": entry.name,
                    "path": str(entry),
                    "type": "folder" if entry.is_dir() else entry.suffix.lower().lstrip("."),
                    "size_kb": round(stat.st_size / 1024, 1) if entry.is_file() else None,
                    "modified": stat.st_mtime,
                }
                if entry.is_dir() and current_depth < depth:
                    item["children"] = _list_dir(entry, current_depth + 1)
                items.append(item)
            except OSError:
                continue
        return items

    children = _list_dir(base, 1)
    return json.dumps({
        "path": str(base),
        "name": base.name,
        "count": len(children),
        "items": children,
    }, indent=2)


def _run_read_text(file_path: str, max_chars: int) -> str:
    """Read a plain-text file from SharePoint/OneDrive."""
    try:
        p = _safe_path(file_path)
    except ValueError as e:
        return json.dumps({"error": "Operation failed"})

    if not p.exists():
        return json.dumps({"error": f"File not found: {file_path}"})
    if not p.is_file():
        return json.dumps({"error": f"Not a file: {file_path}"})

    binary_types = {".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".xls", ".pdf"}
    if p.suffix.lower() in binary_types:
        return json.dumps({
            "error": f"Use sp_read_pptx/sp_read_docx/sp_read_xlsx for {p.suffix} files."
        })

    try:
        content = p.read_text(encoding="utf-8", errors="replace")
        truncated = len(content) > max_chars
        return json.dumps({
            "path": str(p),
            "name": p.name,
            "size_kb": round(p.stat().st_size / 1024, 1),
            "truncated": truncated,
            "content": content[:max_chars],
        }, indent=2)
    except Exception as e:
        return json.dumps({"error": f"Could not read file: {e}"})


def _run_read_pptx(file_path: str) -> str:
    """Extract slide text from a .pptx file."""
    try:
        p = _safe_path(file_path)
    except ValueError as e:
        return json.dumps({"error": "Operation failed"})

    if not p.exists():
        return json.dumps({"error": f"File not found: {file_path}"})

    try:
        from pptx import Presentation
    except ImportError:
        return json.dumps({"error": "python-pptx not installed. Run: pip install python-pptx"})

    try:
        prs = Presentation(str(p))
    except Exception as e:
        return json.dumps({"error": f"Could not open presentation: {e}"})

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        texts.append(line.strip())

        notes_text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            notes_text = " ".join(par.text for par in tf.paragraphs if par.text.strip())

        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        slides.append({
            "slide": i,
            "title": title,
            "content": texts,
            "notes": notes_text,
        })

    return json.dumps({
        "path": str(p),
        "name": p.name,
        "slide_count": len(slides),
        "slides": slides,
    }, indent=2)


def _run_read_docx(file_path: str, max_chars: int) -> str:
    """Extract text from a .docx file."""
    try:
        p = _safe_path(file_path)
    except ValueError as e:
        return json.dumps({"error": "Operation failed"})

    if not p.exists():
        return json.dumps({"error": f"File not found: {file_path}"})

    try:
        import docx
    except ImportError:
        return json.dumps({"error": "python-docx not installed. Run: pip install python-docx"})

    try:
        doc = docx.Document(str(p))
    except Exception as e:
        return json.dumps({"error": f"Could not open document: {e}"})

    paragraphs = []
    total_chars = 0
    truncated = False

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if total_chars + len(text) > max_chars:
            truncated = True
            break
        style = para.style.name if para.style else ""
        paragraphs.append({"style": style, "text": text})
        total_chars += len(text)

    full_text = "\n".join(par["text"] for par in paragraphs)

    return json.dumps({
        "path": str(p),
        "name": p.name,
        "paragraph_count": len(paragraphs),
        "truncated": truncated,
        "text": full_text,
        "paragraphs": paragraphs,
    }, indent=2)


def _run_read_xlsx(file_path: str, sheet_name: str, max_rows: int) -> str:
    """Read data from a .xlsx file."""
    try:
        p = _safe_path(file_path)
    except ValueError as e:
        return json.dumps({"error": "Operation failed"})

    if not p.exists():
        return json.dumps({"error": f"File not found: {file_path}"})

    try:
        import openpyxl
    except ImportError:
        return json.dumps({"error": "openpyxl not installed. Run: pip install openpyxl"})

    try:
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    except Exception as e:
        return json.dumps({"error": f"Could not open workbook: {e}"})

    sheet_names = wb.sheetnames
    sheets_to_read = [sheet_name] if sheet_name else sheet_names

    result_sheets = {}
    for sname in sheets_to_read:
        if sname not in sheet_names:
            result_sheets[sname] = {"error": f"Sheet '{sname}' not found"}
            continue
        ws = wb[sname]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            row_vals = [str(cell) if cell is not None else "" for cell in row]
            if any(v for v in row_vals):
                rows.append(row_vals)
        result_sheets[sname] = {"rows": rows, "row_count": len(rows)}

    wb.close()

    return json.dumps({
        "path": str(p),
        "name": p.name,
        "sheets": sheet_names,
        "data": result_sheets,
    }, indent=2)


def _run_health() -> str:
    """Check health and list available roots."""
    roots_status = []
    for root in ALLOWED_ROOTS:
        try:
            items = list(root.iterdir())
            file_count = sum(1 for i in root.rglob("*") if i.is_file() and not i.name.startswith("."))
            roots_status.append({
                "name": root.name,
                "path": str(root),
                "exists": True,
                "top_level_items": len(items),
                "total_files": file_count,
            })
        except Exception as e:
            roots_status.append({
                "name": root.name,
                "path": str(root),
                "exists": root.exists(),
                "error": "Operation failed",
            })

    return json.dumps({
        "status": "ok",
        "mode": "READ-ONLY",
        "write_protected": True,
        "allowed_roots": len(ALLOWED_ROOTS),
        "roots": roots_status,
    }, indent=2)


# ---------------------------------------------------------------------------
# BaseTool implementation
# ---------------------------------------------------------------------------

def _node_bin() -> str:
    for n in ("/opt/homebrew/bin/node", "/usr/local/bin/node",
              "/opt/homebrew/opt/node@22/bin/node", str(Path.home() / ".npm-global/bin/node")):
        if os.path.exists(n):
            return n
    return "node"


def _run_tenant_search(query: str, max_results: int) -> str:
    """Tenant-wide SharePoint search — everything the user can access across ALL sites (not just
    locally-synced libraries), via the SharePoint Search REST API using his own signed-in session.
    READ-ONLY (search GET only). Returns titles, authors, dates, snippets, and links."""
    max_results = max(1, min(int(max_results or 20), 50))
    script = str(Path.home() / "claude-stack" / "tools" / "sharepoint_search.js")
    try:
        # 30s cap: ChatGPT Actions time out ~45s, so a 90s search surfaced to the user as a
        # connector error even when the backend eventually succeeded. Fail fast instead.
        out = subprocess.run([_node_bin(), script, query, "--rows", str(max_results)],
                             capture_output=True, text=True, timeout=30)
        text = (out.stdout or "").strip()
        if not text:
            return json.dumps({"error": "no results or SharePoint session needs re-login",
                               "detail": (out.stderr or "").strip()[:300]})
        return text[:50000]
    except subprocess.TimeoutExpired:
        return json.dumps({"error": "tenant SharePoint search timed out"})
    except Exception as e:
        return json.dumps({"error": f"tenant SharePoint search failed: {e}"})


class SharePointLocalTool(BaseTool):
    name = "sharepoint_local"
    description = "SharePoint / OneDrive — read-only file search, listing, and reading from synced CloudStorage libraries"

    def get_definitions(self) -> List[dict]:
        return [
            make_tool_def(
                "sp_search",
                "Search for files across all synced SharePoint and OneDrive libraries using macOS Spotlight.",
                {
                    "query": {"type": "string", "description": "Search query — filename keyword, content keyword, or file extension"},
                    "file_types": {"type": "string", "description": "Comma-separated extensions to filter, e.g. 'pptx,docx,xlsx'. Leave empty for all."},
                    "root": {"type": "string", "description": "Restrict search to a specific subfolder path (must be inside CloudStorage). Leave empty for all."},
                    "max_results": {"type": "integer", "description": "Max number of results (default 50, max 200)"},
                },
                ["query"],
            ),
            make_tool_def(
                "sp_list",
                "List files and folders in a SharePoint/OneDrive directory. Leave folder_path empty to list all synced library roots.",
                {
                    "folder_path": {"type": "string", "description": "Path to list. Leave empty to list all synced library roots."},
                    "depth": {"type": "integer", "description": "How many levels to recurse (1 = immediate children, max 3)"},
                },
                [],
            ),
            make_tool_def(
                "sp_read_text",
                "Read a plain-text file (txt, md, csv, json, py, etc.) from SharePoint/OneDrive.",
                {
                    "file_path": {"type": "string", "description": "Absolute path to the file"},
                    "max_chars": {"type": "integer", "description": "Maximum characters to return (default 50000)"},
                },
                ["file_path"],
            ),
            make_tool_def(
                "sp_read_pptx",
                "Extract all slide text and notes from a PowerPoint (.pptx) file in SharePoint/OneDrive.",
                {
                    "file_path": {"type": "string", "description": "Absolute path to the .pptx file"},
                },
                ["file_path"],
            ),
            make_tool_def(
                "sp_read_docx",
                "Extract full text from a Word document (.docx) in SharePoint/OneDrive.",
                {
                    "file_path": {"type": "string", "description": "Absolute path to the .docx file"},
                    "max_chars": {"type": "integer", "description": "Maximum characters to return (default 100000)"},
                },
                ["file_path"],
            ),
            make_tool_def(
                "sp_read_xlsx",
                "Read data from an Excel spreadsheet (.xlsx) in SharePoint/OneDrive.",
                {
                    "file_path": {"type": "string", "description": "Absolute path to the .xlsx file"},
                    "sheet_name": {"type": "string", "description": "Name of the sheet to read. Leave empty for all sheets."},
                    "max_rows": {"type": "integer", "description": "Maximum rows per sheet (default 500)"},
                },
                ["file_path"],
            ),
            make_tool_def(
                "sp_search_all",
                "Tenant-wide SharePoint search across ALL company sites the user can access (not just "
                "synced libraries) — finds tribal-knowledge docs, other teams' files, site pages, and "
                "shared docs. Use this for broad 'find/search in SharePoint' or company-context questions. "
                "Read-only. Returns titles, authors, dates, content snippets, and links.",
                {
                    "query": {"type": "string", "description": "Search terms (keywords, ticket #, topic, person)."},
                    "max_results": {"type": "integer", "description": "Max results (default 20, max 50)."},
                },
                ["query"],
            ),
            make_tool_def(
                "sp_health",
                "Check SharePoint/OneDrive health and list available synced library roots with file counts.",
                {},
                [],
            ),
        ]

    async def handle(self, tool_name: str, tool_input: dict) -> str:
        loop = asyncio.get_event_loop()

        if tool_name == "sp_search":
            query = tool_input["query"]
            file_types = tool_input.get("file_types", "")
            root = tool_input.get("root", "")
            max_results = tool_input.get("max_results", 50)
            return await loop.run_in_executor(None, _run_search, query, file_types, root, max_results)

        elif tool_name == "sp_list":
            folder_path = tool_input.get("folder_path", "")
            depth = tool_input.get("depth", 1)
            return await loop.run_in_executor(None, _run_list, folder_path, depth)

        elif tool_name == "sp_read_text":
            file_path = tool_input["file_path"]
            max_chars = tool_input.get("max_chars", 50000)
            return await loop.run_in_executor(None, _run_read_text, file_path, max_chars)

        elif tool_name == "sp_read_pptx":
            file_path = tool_input["file_path"]
            return await loop.run_in_executor(None, _run_read_pptx, file_path)

        elif tool_name == "sp_read_docx":
            file_path = tool_input["file_path"]
            max_chars = tool_input.get("max_chars", 100000)
            return await loop.run_in_executor(None, _run_read_docx, file_path, max_chars)

        elif tool_name == "sp_read_xlsx":
            file_path = tool_input["file_path"]
            sheet_name = tool_input.get("sheet_name", "")
            max_rows = tool_input.get("max_rows", 500)
            return await loop.run_in_executor(None, _run_read_xlsx, file_path, sheet_name, max_rows)

        elif tool_name == "sp_search_all":
            query = tool_input["query"]
            max_results = tool_input.get("max_results", 20)
            return await loop.run_in_executor(None, _run_tenant_search, query, max_results)

        elif tool_name == "sp_health":
            return await loop.run_in_executor(None, _run_health)

        else:
            return f"Unknown sharepoint tool: {tool_name}"
