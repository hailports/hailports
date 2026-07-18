#!/usr/bin/env python3
"""Stdlib-light text extraction from common attachment types. No external services."""
from __future__ import annotations

import email
import html
import io
import os
import re
import subprocess
import zipfile


_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic"}


def _vision_ocr(path: str, timeout: int = 35) -> str:
    """Use the existing local macOS Vision helper when available; no network or cloud API."""
    helper = os.path.expanduser(os.environ.get("ATTACHMENT_VISION_OCR_BIN", "~/.zoom-mcp-ocr"))
    if not os.path.isfile(helper) or not os.access(helper, os.X_OK):
        return ""
    try:
        result = subprocess.run([helper, path], capture_output=True, text=True, timeout=timeout)
        text = (result.stdout or "").strip()
        if result.returncode == 0 and text and not text.startswith("ERROR|"):
            return text
    except Exception:
        pass
    return ""


def _ocr(source, timeout: int = 35) -> str:
    """Local-only OCR. `source` is a path or PNG/JPEG bytes for stdin."""
    cmd = ["tesseract", "stdin" if isinstance(source, bytes) else str(source), "stdout", "--psm", "6"]
    try:
        result = subprocess.run(
            cmd,
            input=source if isinstance(source, bytes) else None,
            capture_output=True,
            timeout=timeout,
        )
        if result.returncode == 0:
            return result.stdout.decode("utf-8", "ignore").strip()
    except Exception:
        pass
    return ""


def _image(path: str) -> str:
    # Vision's accurate recognizer handles dense UI screenshots materially better than Tesseract.
    # Tesseract remains the portable fallback for hosts without the already-built local helper.
    text = _vision_ocr(path) or _ocr(path)
    if text:
        return text
    # Tesseract cannot decode HEIC on every host. Normalize through Pillow when available.
    try:
        from PIL import Image
        image = Image.open(path)
        image.thumbnail((5000, 5000))
        payload = io.BytesIO()
        image.convert("RGB").save(payload, format="PNG")
        text = _ocr(payload.getvalue())
    except Exception:
        text = ""
    return text or "[Image present; local OCR returned no readable text]"


def _eml(path: str) -> str:
    m = email.message_from_bytes(open(path, "rb").read())
    out = [f"Subject: {m.get('Subject', '')}", f"From: {m.get('From', '')}", f"Date: {m.get('Date', '')}", ""]
    for p in m.walk():
        if p.get_content_type() == "text/plain":
            try:
                out.append((p.get_payload(decode=True) or b"").decode("utf-8", "ignore"))
            except Exception:
                pass
    return "\n".join(out)


def _pdf(path: str) -> str:
    try:
        r = subprocess.run(["pdftotext", "-q", path, "-"], capture_output=True, timeout=25)
        if r.returncode == 0 and r.stdout.strip():
            return r.stdout.decode("utf-8", "ignore")
    except Exception:
        pass
    try:
        import fitz  # type: ignore
        document = fitz.open(path)
        text = "\n".join(page.get_text("text") for page in document).strip()
        if text:
            document.close()
            return text
        # Scanned PDFs have no text layer. OCR a bounded number of pages locally.
        pages = []
        for page in list(document)[:8]:
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            page_text = _ocr(pixmap.tobytes("png"), timeout=35)
            if page_text:
                pages.append(page_text)
        document.close()
        return "\n\n".join(pages) or "[PDF present; local extraction/OCR returned no readable text]"
    except Exception:
        try:
            from pdfminer.high_level import extract_text as _et  # type: ignore
            return _et(path) or "[PDF present; local extraction returned no readable text]"
        except Exception:
            return "[PDF present; no local text extractor available]"


def _docx(path: str) -> str:
    try:
        from docx import Document  # type: ignore
        document = Document(path)
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    parts.append("\t".join(values))
        return "\n".join(parts)
    except Exception:
        try:
            xml = zipfile.ZipFile(path).read("word/document.xml").decode("utf-8", "ignore")
            xml = re.sub(r"</w:p>", "\n", xml)
            return html.unescape(re.sub(r"[ \t]+", " ", re.sub(r"<[^>]+>", " ", xml))).strip()
        except Exception as e:
            return f"[docx parse error: {e}]"


def _xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
        workbook = load_workbook(path, read_only=True, data_only=False)
        parts: list[str] = []
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_number > 2_000:
                    parts.append("[remaining rows omitted]")
                    break
                values = [str(value).strip() if value is not None else "" for value in row[:100]]
                while values and not values[-1]:
                    values.pop()
                if any(values):
                    parts.append("\t".join(values))
        workbook.close()
        return "\n".join(parts)
    except Exception:
        try:
            z = zipfile.ZipFile(path)
            parts = []
            if "xl/sharedStrings.xml" in z.namelist():
                xml = z.read("xl/sharedStrings.xml").decode("utf-8", "ignore")
                parts += re.findall(r"<t[^>]*>([^<]+)</t>", xml)
            return " | ".join(html.unescape(value) for value in parts)
        except Exception as e:
            return f"[xlsx parse error: {e}]"


def _pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
        presentation = Presentation(path)
        parts: list[str] = []
        for number, slide in enumerate(presentation.slides, start=1):
            values = [shape.text.strip() for shape in slide.shapes
                      if hasattr(shape, "text") and shape.text.strip()]
            if values:
                parts.append(f"[Slide {number}]\n" + "\n".join(values))
        return "\n\n".join(parts)
    except Exception:
        try:
            z = zipfile.ZipFile(path)
            parts = []
            for n in z.namelist():
                if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                    parts += re.findall(r"<a:t>([^<]+)</a:t>", z.read(n).decode("utf-8", "ignore"))
            return " | ".join(html.unescape(value) for value in parts)
        except Exception as e:
            return f"[pptx parse error: {e}]"


def extract_text(path: str, cap: int = 50000) -> str:
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".eml":
            t = _eml(path)
        elif ext == ".pdf":
            t = _pdf(path)
        elif ext == ".docx":
            t = _docx(path)
        elif ext == ".xlsx":
            t = _xlsx(path)
        elif ext == ".pptx":
            t = _pptx(path)
        elif ext in _IMAGE_EXTS:
            t = _image(path)
        elif ext in (".txt", ".csv", ".md", ".log", ".json", ".eml", ".rtf", ".html",
                     ".htm", ".xml", ".yaml", ".yml", ".sql"):
            t = open(path, encoding="utf-8", errors="ignore").read()
        else:
            raw = open(path, "rb").read(2_000_000)
            t = "\n".join(re.findall(r"[\x20-\x7e]{6,}", raw.decode("latin-1", "ignore")))
        return (t or "").strip()[:cap]
    except Exception as e:
        return f"[extract error: {e}]"
